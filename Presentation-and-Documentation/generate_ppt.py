import os

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()


def add_bullet_slide(title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    body = s.shapes.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = text
        p.level = 0
        p.font.size = Pt(18)


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "House Price Prediction — Task 4"
slide.placeholders[1].text = "End-to-end Data Science Project (Bangalore-focused demo)"

# Slide 2: Problem Statement
add_bullet_slide("Problem Statement", [
    "Goal: predict house sale price from property and location features.",
    "Use case: quick pricing estimate for buyers, sellers, and agents.",
    "Inputs: square_feet, bedrooms, bathrooms, age, garage, location_score, condition.",
])

# Slide 3: Data Pipeline
add_bullet_slide("Data & Preprocessing Pipeline", [
    "Synthetic housing dataset generated for fast model prototyping.",
    "AdvancedPreprocessor handles missing values, condition encoding, scaling, and outlier handling.",
    "Same preprocessor is saved and reused during API inference for consistency.",
])

# Slide 4: Model Training
add_bullet_slide("Model Training & Selection", [
    "Trained XGBoost, LightGBM, CatBoost, RandomForest, and GradientBoosting.",
    "Compared model performance and selected best_model.pkl for deployment.",
    "All artifacts are versioned in models/ for reproducible predictions.",
])

# Slide 5: API and Frontend
add_bullet_slide("House Price API + Frontend", [
    "FastAPI endpoints: /predict, /batch-predict, /health, /model-info, /feature-importance.",
    "Frontend supports single prediction, CSV batch upload, and analytics view.",
    "INR display toggle added for Bangalore presentation context.",
])

# Slide 6: Example Prediction Flow
add_bullet_slide("Example Prediction Flow", [
    "Example input: 1500 sq ft, 3 bed, 2 bath, age 10, garage 1, location_score 8, condition Good.",
    "API response returns predicted_price with confidence and selected model name.",
    "Prediction is bounded to practical min/max price ranges for stable output.",
])

# Slide 7: Model Insight Screenshot
slide7 = prs.slides.add_slide(prs.slide_layouts[5])
slide7.shapes.title.text = "Model Insight: Feature Importance"
img_path = "Presentation-and-Documentation/model_importance.png"
if os.path.exists(img_path):
    slide7.shapes.add_picture(img_path, Inches(0.8), Inches(1.4), width=Inches(8.0))
else:
    tx = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.0), Inches(1.5)).text_frame
    tx.text = "model_importance.png not found. Run create_model_screenshot.py first."
    tx.paragraphs[0].font.size = Pt(18)

# Slide 8: Deployment + Next Steps
add_bullet_slide("Deployment & Next Steps", [
    "Docker and docs prepared for portable deployment and demo.",
    "Next: train on real Bangalore housing data for production-level accuracy.",
    "Add CI tests, model monitoring, and automated retraining pipeline.",
])

out_path = "Presentation-and-Documentation/Task-4-Presentation.pptx"
prs.save(out_path)
print("saved", out_path)
