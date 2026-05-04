from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

PRIMARY_COLOR = RGBColor(52, 152, 219)
SECONDARY_COLOR = RGBColor(46, 62, 80)
WHITE = RGBColor(255, 255, 255)
TEXT_COLOR = RGBColor(52, 73, 94)
LIGHT_BG = RGBColor(243, 246, 249)
CARD_BG = RGBColor(255, 255, 255)


def add_slide_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_header(slide, title, subtitle=None):
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.color.rgb = PRIMARY_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.14), Inches(6.8), Inches(0.5))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(6.8), Inches(0.18), Inches(2.8), Inches(0.4))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        p = sub_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.RIGHT


def style_box(text_frame, font_size=15, bold_prefix=False):
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = Pt(8)
    text_frame.margin_right = Pt(8)
    text_frame.margin_top = Pt(4)
    text_frame.margin_bottom = Pt(4)
    for paragraph in text_frame.paragraphs:
        paragraph.space_after = Pt(3)
        paragraph.space_before = Pt(0)
        paragraph.line_spacing = 1.0
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = TEXT_COLOR


def add_text_card(slide, left, top, width, height, lines, title=None, font_size=15):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(220, 228, 236)

    box_top = top + Inches(0.12) if title else top + Inches(0.05)
    box_height = height - Inches(0.12) if title else height - Inches(0.08)
    box = slide.shapes.add_textbox(left + Inches(0.12), box_top, width - Inches(0.24), box_height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if title:
        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(6)

    for index, line in enumerate(lines):
        p = frame.add_paragraph() if (title or index > 0) else frame.paragraphs[0]
        if not title and index == 0:
            pass
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(2)
        if line.startswith("•"):
            p.level = 1

    style_box(frame, font_size=font_size)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PRIMARY_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.8))
    sub_frame = subtitle_box.text_frame
    sub_frame.clear()
    p = sub_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_two_column_slide(title, left_title, left_lines, right_title, right_lines, subtitle="CodTech IT Solutions"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_header(slide, title, subtitle)
    add_text_card(slide, Inches(0.45), Inches(1.2), Inches(4.45), Inches(5.9), left_lines, left_title, 14)
    add_text_card(slide, Inches(5.1), Inches(1.2), Inches(4.45), Inches(5.9), right_lines, right_title, 14)


def add_single_card_slide(title, lines, subtitle="CodTech IT Solutions"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_header(slide, title, subtitle)
    add_text_card(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.9), lines, None, 15)


def add_image_slide(title, image_path, caption, subtitle="CodTech IT Solutions"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide)
    add_header(slide, title, subtitle)

    picture_left = Inches(0.55)
    picture_top = Inches(1.25)
    picture_width = Inches(8.9)
    picture_height = Inches(5.35)
    slide.shapes.add_picture(str(image_path), picture_left, picture_top, width=picture_width, height=picture_height)

    caption_box = slide.shapes.add_textbox(Inches(0.65), Inches(6.75), Inches(8.7), Inches(0.35))
    frame = caption_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = caption
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER


add_title_slide("INTERNSHIP PROJECT PRESENTATION", "CodTech IT Solutions")

company_left = [
    "• 8+ years of experience",
    "• 1500+ projects completed",
    "• 1000+ happy clients",
    "• 750+ skilled team members",
]
company_right = [
    "• Google, Microsoft, and AWS partnerships",
    "• MSME, AICTE, MCA, ISO 9001:2015 certifications",
    "• Services: Web, mobile, marketing, UI/UX",
    "• Focus: professional digital solutions that drive growth",
]
add_two_column_slide("COMPANY INTRODUCTION", "Company Snapshot", company_left, "Key Highlights", company_right)

task1_left = [
    "• Built an ETL pipeline in Python",
    "• Removed duplicates and handled missing values",
    "• Scaled numeric columns",
]
task1_right = [
    "• Encoded categorical columns",
    "• Saved transformed output as CSV",
    "• Tools: pandas, scikit-learn, StandardScaler, SimpleImputer, OneHotEncoder",
]
add_two_column_slide("TASK 1 - ETL PIPELINE DEVELOPMENT", "Objective", task1_left, "Implementation", task1_right)

task1_screenshot = Path(r"E:\ProjectsGithub\Data Science Internship\task1_etl_screenshot.png")
if task1_screenshot.exists():
    add_image_slide(
        "TASK 1 SCREENSHOT",
        task1_screenshot,
        "ETL pipeline code screenshot showing preprocessing, transformation, and loading steps.",
    )

task2_left = [
    "• CNN-based digit classification on MNIST",
    "• Conv2D + MaxPooling layers",
    "• Dense layer with dropout",
]
task2_right = [
    "• Generated training curves",
    "• Created confusion matrix output",
    "• Tools: TensorFlow/Keras, NumPy, Matplotlib, scikit-learn",
]
add_two_column_slide("TASK 2 - DEEP LEARNING PROJECT", "Model", task2_left, "Outputs", task2_right)

task2_screenshot = Path(r"E:\ProjectsGithub\Data Science Internship\task2_deep_learning_screenshot.png")
if task2_screenshot.exists():
    add_image_slide(
        "TASK 2 SCREENSHOT",
        task2_screenshot,
        "Deep learning code screenshot showing the CNN model architecture and evaluation functions.",
    )

task3_left = [
    "• Interactive product gallery with real images",
    "• Thumbnail navigation and image counter",
    "• Customer reviews and rating summary",
]
task3_right = [
    "• Add to cart, wishlist, and review form",
    "• Responsive layout for desktop and mobile",
    "• Tools: HTML5, CSS3, JavaScript, Unsplash images",
]
add_two_column_slide("TASK 3 - E-COMMERCE PRODUCT PAGE", "Features", task3_left, "Tech Stack", task3_right)

task3_screenshot = Path(r"E:\ProjectsGithub\Data Science Internship\task3_product_page_screenshot.png")
if task3_screenshot.exists():
    add_image_slide(
        "TASK 3 SCREENSHOT",
        task3_screenshot,
        "Interactive product page screenshot showing the gallery, product details, and review section.",
    )

tools_left = [
    "• Python",
    "• pandas, NumPy",
    "• scikit-learn",
    "• TensorFlow/Keras",
]
tools_right = [
    "• HTML5, CSS3, JavaScript",
    "• VS Code",
    "• Git and GitHub",
    "• PowerPoint via python-pptx",
]
add_two_column_slide("TOOLS & TECHNOLOGIES USED", "Programming", tools_left, "Development Tools", tools_right)

summary_lines = [
    "• Task 1: ETL pipeline development",
    "• Task 2: Deep learning image classification",
    "• Task 3: E-commerce product page design",
    "• Learned data processing, ML, and frontend development",
    "• Internship work completed for CodTech IT Solutions",
]
add_single_card_slide("PROJECT SUMMARY", summary_lines)

output_path = r"E:\ProjectsGithub\Data Science Internship\CODTECH_INTERNSHIP_PRESENTATION_FINAL_WITH_TASK1_TASK2_SS.pptx"
prs.save(output_path)
print("Presentation rebuilt successfully")
print(output_path)
