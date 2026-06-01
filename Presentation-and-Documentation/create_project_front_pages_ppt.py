from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Project Front Pages — 2025-26 MCA"
subtitle = slide.placeholders[1]
subtitle.text = "Full Data Science Project: End-to-end pipeline, API & Frontend Demo"

# Info slide
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Project Details"
body = s.shapes.placeholders[1].text_frame
body.clear()
items = [
    "Project: Full DS Project (Task-4)",
    "Author: Internship Team — 2025-26 MCA cohort",
    "Scope: Data preprocessing → Model training → FastAPI deployment → Frontend UI",
    "Artifacts: serialized preprocessor & models (models/), API (src/api.py), frontend (/ui)",
    "Repo: branch 'task-4' with Docker & docs for reproducibility",
]
for i, t in enumerate(items):
    p = body.add_paragraph() if i>0 else body.paragraphs[0]
    p.text = t
    p.level = 0
    p.font.size = Pt(18)

# Summary slide
s2 = prs.slides.add_slide(prs.slide_layouts[1])
s2.shapes.title.text = "What We Delivered"
body2 = s2.shapes.placeholders[1].text_frame
body2.clear()
points = [
    "Synthetic data pipeline and AdvancedPreprocessor",
    "Trained multiple models; saved best-performing artifacts",
    "FastAPI endpoints for single and batch predictions",
    "Interactive frontend with INR display and CSV upload",
    "Presentation and documentation (this PPT, READMEs, deployment guides)",
]
for i, t in enumerate(points):
    p = body2.add_paragraph() if i>0 else body2.paragraphs[0]
    p.text = t
    p.level = 0
    p.font.size = Pt(18)

# Footer slide with date
s3 = prs.slides.add_slide(prs.slide_layouts[1])
s3.shapes.title.text = "Contact & Next Steps"
body3 = s3.shapes.placeholders[1].text_frame
body3.clear()
notes = [
    "Next: Train on real Bangalore dataset and CI integration",
    "Push final assets to 'task-4' branch and open PR for review",
    f"Generated: {datetime.now().strftime('%Y-%m-%d')}",
]
for i, t in enumerate(notes):
    p = body3.add_paragraph() if i>0 else body3.paragraphs[0]
    p.text = t
    p.level = 0
    p.font.size = Pt(18)

# Save with same folder and matching name
out_path = 'Presentation-and-Documentation/project front pages 2025-26-MCA.pptx'
prs.save(out_path)
print('saved', out_path)
